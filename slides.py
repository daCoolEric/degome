"""
Course-material text extraction for study guides.

Native (pure-Python, always available):
    .pdf   -> pypdf
    .pptx  -> python-pptx (slides, tables, speaker notes)
    .docx  -> python-docx (paragraphs, tables)
    .xlsx  -> openpyxl (cell values per sheet)

Images (.png .jpg .jpeg) -> OCR via tesseract, when installed.

Legacy Office (.ppt .doc .xls) -> converted to the modern format with
LibreOffice when installed, then extracted natively. Without LibreOffice
the user gets a clear "save as .pptx/.docx/.xlsx" message.
"""

import shutil
import subprocess
import tempfile
from pathlib import Path

MAX_CHARS = 400_000  # safety cap for gigantic decks/workbooks

NATIVE = {".pdf", ".pptx", ".docx", ".xlsx"}
LEGACY = {".ppt": ".pptx", ".doc": ".docx", ".xls": ".xlsx"}
IMAGES = {".png", ".jpg", ".jpeg"}
SUPPORTED = NATIVE | set(LEGACY) | IMAGES


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix in LEGACY:
        text = _extract_legacy(path, LEGACY[suffix])
    elif suffix == ".pdf":
        text = _extract_pdf(path)
    elif suffix == ".pptx":
        text = _extract_pptx(path)
    elif suffix == ".docx":
        text = _extract_docx(path)
    elif suffix == ".xlsx":
        text = _extract_xlsx(path)
    elif suffix in IMAGES:
        text = _extract_image(path)
    else:
        raise RuntimeError(
            "Unsupported file type. Accepted: PDF, PPTX/PPT, DOCX/DOC, "
            "XLSX/XLS, PNG, JPG.")

    text = text.strip()
    if len(text) < 30:
        raise RuntimeError(
            "Almost no text could be extracted from this file. If it's a "
            "scanned document, try a clearer photo or a text-based export.")
    return text[:MAX_CHARS]


# ---------------------------------------------------------------- native
def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        content = (page.extract_text() or "").strip()
        if content:
            pages.append(f"--- Slide/Page {i} ---\n{content}")
    return "\n\n".join(pages)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        chunks.append(line)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        chunks.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(f"[Speaker notes] {notes}")
        if chunks:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(chunks))
    return "\n\n".join(slides)


def _extract_docx(path: Path) -> str:
    import docx as docx_lib

    doc = docx_lib.Document(str(path))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r > 2000:  # sanity cap per sheet
                rows.append("... (sheet truncated)")
                break
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"--- Sheet: {ws.title} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


# ---------------------------------------------------------------- images
def _extract_image(path: Path) -> str:
    if shutil.which("tesseract") is None:
        raise RuntimeError(
            "Image OCR needs tesseract, which isn't installed on this server. "
            "Install it (apt install tesseract-ocr) or upload a PDF/PPTX instead.")
    import pytesseract
    from PIL import Image, ImageOps

    img = Image.open(str(path))
    img = ImageOps.exif_transpose(img)  # respect phone rotation
    if img.mode != "L":
        img = img.convert("L")  # grayscale helps OCR on photos
    text = pytesseract.image_to_string(img)
    return f"--- Image: {path.name} (OCR) ---\n{text}"


# ---------------------------------------------------------------- legacy
def _find_soffice() -> str | None:
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    # Windows installs aren't on PATH by default
    import os
    for base in (os.environ.get("ProgramFiles", r"C:\Program Files"),
                 os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)")):
        candidate = Path(base) / "LibreOffice" / "program" / "soffice.exe"
        if candidate.exists():
            return str(candidate)
    return None


def _extract_legacy(path: Path, modern_suffix: str) -> str:
    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError(
            f"Legacy {path.suffix} files need LibreOffice installed on the "
            f"machine running Degome (free, libreoffice.org). Or open the file "
            f"in PowerPoint/Word/Excel and save it as {modern_suffix}, then "
            f"upload that.")
    with tempfile.TemporaryDirectory() as tmp:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", modern_suffix.lstrip("."),
             "--outdir", tmp, str(path)],
            capture_output=True, text=True, timeout=180,
        )
        converted = Path(tmp) / (path.stem + modern_suffix)
        if result.returncode != 0 or not converted.exists():
            raise RuntimeError(
                f"Could not convert this {path.suffix} file. Open it and save "
                f"as {modern_suffix}, then upload that.")
        return extract_text(converted)
